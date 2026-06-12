from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from utils import logger
from config import (
    OUTPUT_PPT, PPT_SLIDE_WIDTH, PPT_SLIDE_HEIGHT,
    IMAGE_MAX_WIDTH, IMAGE_MAX_HEIGHT
)
from PIL import Image

class PPTGenerator:
    """PPT生成器"""

    def __init__(self, output_path=OUTPUT_PPT):
        self.output_path = output_path
        self.prs = Presentation()
        # 设置幻灯片尺寸
        self.prs.slide_width = Inches(PPT_SLIDE_WIDTH)
        self.prs.slide_height = Inches(PPT_SLIDE_HEIGHT)
        logger.info("PPT生成器初始化完成")

    def add_slide(self, data):
        """
        添加一页幻灯片

        Args:
            data: dict {
                'media_name': 媒体名称,
                'title': 标题,
                'publish_time': 发布时间,
                'link': 链接,
                'screenshot': 截图路径,
                'summary': 摘要文本
            }
        """
        try:
            # 使用空白布局
            blank_slide_layout = self.prs.slide_layouts[6]
            slide = self.prs.slides.add_slide(blank_slide_layout)

            # 标题（顶部）
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(9), Inches(0.7)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            p = title_frame.paragraphs[0]
            p.text = data.get('title', '')
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)

            # 媒体名称和发布时间（标题下方同一行）
            media_time_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.05),
                Inches(9), Inches(0.35)
            )
            media_time_frame = media_time_box.text_frame
            p = media_time_frame.paragraphs[0]
            p.text = f"{data.get('media_name', '')}    发布时间: {data.get('publish_time', '')}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(100, 100, 100)

            # 链接（媒体信息下方）
            link_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.45),
                Inches(9), Inches(0.3)
            )
            link_frame = link_box.text_frame
            link_frame.word_wrap = True
            p = link_frame.paragraphs[0]
            p.text = f"链接: {data.get('link', '')}"
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(0, 102, 204)  # 蓝色

            # 图片区域（中间）
            screenshot_path = data.get('screenshot')
            image_top = Inches(1.9)  # 调整位置

            if screenshot_path and os.path.exists(screenshot_path):
                # 计算图片尺寸，保持比例
                left, top, width, height = self._calculate_image_size(
                    screenshot_path,
                    max_width=IMAGE_MAX_WIDTH,
                    max_height=IMAGE_MAX_HEIGHT,
                    top_position=image_top
                )
                slide.shapes.add_picture(screenshot_path, left, top, width, height)
                logger.info(f"添加图片: {screenshot_path}")

                # 如果有摘要，放在图片下方
                summary = data.get('summary', '')
                if summary:
                    summary_top = top + height + Inches(0.2)
                    summary_height = Inches(PPT_SLIDE_HEIGHT) - summary_top - Inches(0.3)

                    if summary_height > Inches(0.5):  # 确保有足够空间
                        summary_box = slide.shapes.add_textbox(
                            Inches(0.5), summary_top,
                            Inches(9), summary_height
                        )
                        summary_frame = summary_box.text_frame
                        summary_frame.word_wrap = True
                        p = summary_frame.paragraphs[0]
                        p.text = summary
                        p.font.size = Pt(11)
                        p.font.color.rgb = RGBColor(50, 50, 50)
                        p.alignment = PP_ALIGN.LEFT
            else:
                # 没有图片时，只显示摘要或链接
                summary = data.get('summary', '') or f"链接: {data.get('link', '')}"
                summary_box = slide.shapes.add_textbox(
                    Inches(0.5), image_top,
                    Inches(9), Inches(4)
                )
                summary_frame = summary_box.text_frame
                summary_frame.word_wrap = True
                p = summary_frame.paragraphs[0]
                p.text = summary
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(50, 50, 50)

            logger.info(f"成功添加幻灯片: {data.get('title', '')[:30]}...")

        except Exception as e:
            logger.error(f"添加幻灯片失败: {e}")
            raise

    def _calculate_image_size(self, image_path, max_width, max_height, top_position):
        """
        计算图片显示尺寸，保持原始比例

        Returns:
            tuple: (left, top, width, height) in Inches
        """
        try:
            with Image.open(image_path) as img:
                img_width, img_height = img.size

            # 计算缩放比例
            width_ratio = max_width / (img_width / 96)  # 96 DPI
            height_ratio = max_height / (img_height / 96)
            scale = min(width_ratio, height_ratio, 1.0)  # 不放大

            # 计算最终尺寸
            final_width = Inches((img_width / 96) * scale)
            final_height = Inches((img_height / 96) * scale)

            # 居中显示
            left = Inches((PPT_SLIDE_WIDTH - (img_width / 96) * scale) / 2)
            top = top_position

            return left, top, final_width, final_height

        except Exception as e:
            logger.warning(f"计算图片尺寸失败: {e}, 使用默认尺寸")
            return Inches(1), top_position, Inches(8), Inches(4.5)

    def save(self):
        """保存PPT文件"""
        try:
            # 确保输出目录存在
            output_dir = os.path.dirname(self.output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)

            self.prs.save(self.output_path)
            logger.info(f"PPT保存成功: {self.output_path}")
            return self.output_path

        except Exception as e:
            logger.error(f"保存PPT失败: {e}")
            raise

if __name__ == '__main__':
    # 测试代码
    generator = PPTGenerator('test_output.pptx')

    test_data = {
        'media_name': 'CCTV-13',
        'title': '测试新闻标题',
        'publish_time': '2026-01-01',
        'link': 'https://example.com',
        'screenshot': None,
        'summary': '这是一段测试摘要文字。' * 20
    }

    generator.add_slide(test_data)
    generator.save()
    print("测试PPT生成完成")
